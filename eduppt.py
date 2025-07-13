import google.generativeai as genai
import requests
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
import os
import random
import re
import string
from googletrans import Translator
from googleapiclient.discovery import build

# === API KEYS ===
GENAI_API_KEY = "AIzaSyDtX3Jdl8Kz_mDyrkdHBTWB3MNNwaXFEPo"
UNSPLASH_ACCESS_KEY = "x7P3BqymG6BFDaFhnAiI_1ROJXsYt3U8xfpoLl5fEuM"
YOUTUBE_API_KEY = "AIzaSyD69g0uo0ldfuCrCKlEzCeMQao8QgrDGDY"

# === Configuration ===
genai.configure(api_key=GENAI_API_KEY)
translator = Translator()

COLOR_PALETTE = [
    (46, 204, 113), (52, 152, 219), (155, 89, 182), (241, 196, 15),
    (230, 126, 34), (231, 76, 60), (149, 165, 166), (52, 73, 94)
]

FONT_MAP = {
    "urdu": "Jameel Noori Nastaleeq",
    "marathi": "Mangal",
    "hindi": "Mangal",
    "english": "Calibri"
}

FONT_LIST = ["Georgia", "Garamond", "Trebuchet MS", "Segoe UI", "Calibri"]


def is_dark_color(rgb_color):
    r, g, b = rgb_color
    brightness = (r * 299 + g * 587 + b * 114) / 1000
    return brightness < 128


def sanitize_filename(filename):
    valid_chars = f"-_.() {string.ascii_letters}{string.digits}"
    return ''.join(c for c in filename if c in valid_chars)


def translate_to_english(text, lang='auto'):
    try:
        translated = translator.translate(text, src=lang, dest='en')
        return translated.text
    except Exception as e:
        print(f"\u274c Translation error: {e}")
        return text


def fetch_unsplash_image(subject, query, language='English'):
    # Remove asterisks and extra whitespace from query
    search_term = f"{subject} {query}".replace('*', '').strip()
    if language.lower() != 'english':
        search_term = translate_to_english(search_term, lang=language.lower())
    search_term = search_term.strip().replace(" ", "+")

    url = f"https://api.unsplash.com/photos/random?query={search_term}&orientation=landscape&client_id={UNSPLASH_ACCESS_KEY}"

    try:
        response = requests.get(url)
        if response.status_code == 200:
            data = response.json()
            image_url = data['urls']['regular']
            img_data = requests.get(image_url).content
            return BytesIO(img_data)
        else:
            print(f"\u26a0\ufe0f Unsplash error: {response.status_code} - {response.text}")
    except Exception as e:
        print(f"\u274c Image fetch error: {e}")
    return None


def search_youtube_video(query):
    try:
        youtube = build("youtube", "v3", developerKey=YOUTUBE_API_KEY)
        request = youtube.search().list(
            part="snippet",
            q=query,
            type="video",
            maxResults=1
        )
        response = request.execute()
        if response["items"]:
            video_id = response["items"][0]["id"]["videoId"]
            return f"https://www.youtube.com/watch?v={video_id}"
    except Exception as e:
        print(f"\u26a0\ufe0f YouTube search error: {e}")
    return None


def split_into_bullets(text):
    lines = text.split('\n')
    bullets = []
    for line in lines:
        line = line.strip()
        if not line:
            continue
        if line.startswith("\u2022") or line.startswith("-"):
            clean = re.sub(r"^[-\u2022\s]+", "", line)
            bullets.append(f"\u2022 {clean}")
        else:
            parts = re.split(r'(?<=[.!?])\s+', line)
            for part in parts:
                part = part.strip()
                if part:
                    bullets.append(f"\u2022 {part}")
    return bullets


def clean_title(title):
    import re
    # Remove all special characters except letters, numbers, and spaces
    title = re.sub(r'[^\w\s]', '', title)
    title = title.strip()
    # Remove all leading/trailing 'u', 'U', 'u ', or ' U' (case-insensitive)
    while title.lower().startswith('u '):
        title = title[2:].strip()
    while title.lower().endswith(' u'):
        title = title[:-2].strip()
    while title.lower().startswith('u') and len(title) > 1:
        title = title[1:].strip()
    while title.lower().endswith('u') and len(title) > 1:
        title = title[:-1].strip()
    # If after cleaning, title is empty, use a placeholder
    if not title:
        title = 'Untitled Slide'
    return title


def generate_ppt_content(class_level, subject, topic, language='English', num_slides=5):
    model = genai.GenerativeModel('gemini-2.0-flash')
    prompt = f"""Create a {num_slides}-slide PowerPoint presentation for a {class_level}th grade {subject} class on the topic \"{topic}\".
Use only formal, culturally accurate, and age-appropriate {language} language.
Do not mix English with {language}. No Roman script or SMS-style writing.
Each slide should include:
- A meaningful, localized title (bold and underlined using *<u>Title</u>* format).
- 5-6 clear bullet points using (\u2022 or -), with examples or facts.
- Do not write paragraphs, just concise bullet points.
- Avoid any unrelated historical figures or general knowledge.
Respond exactly in this format:

SLIDE 1: [Title]
[\u2022 Bullet or - Bullet or sentence]

SLIDE 2: [Title]
[\u2022 Bullet]
[...]"""
    try:
        response = model.generate_content(prompt)
        return format_ppt_output(response.text, subject, topic, language)
    except Exception as e:
        return f"\u274c Error generating content: {str(e)}"


def format_ppt_output(raw_text, subject=None, topic=None, language=None):
    slides = []
    current_slide = {}
    for line in raw_text.split('\n'):
        if line.strip().startswith('SLIDE'):
            if current_slide:
                slides.append(current_slide)
            parts = line.split(':', 1)
            current_slide = {'title': clean_title(parts[1].strip()), 'content': []}
        elif line.strip():
            current_slide['content'].append(line.strip())
    if current_slide:
        slides.append(current_slide)
    # Do NOT add a Practice Questions slide
    return slides


def create_powerpoint(slides, topic, class_level, subject, language="English"):
    prs = Presentation()
    language_key = language.lower().strip()
    rtl = True if language_key == 'urdu' else False

    yt_query = f"{subject} {topic} {language}"
    yt_url = search_youtube_video(yt_query)

    bg_rgb = random.choice(COLOR_PALETTE)
    font_name = FONT_MAP.get(language_key, random.choice(FONT_LIST))
    is_dark = is_dark_color(bg_rgb)
    font_color = RGBColor(255, 255, 255) if is_dark else RGBColor(0, 0, 0)

    title_slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_slide.background.fill.solid()
    title_slide.background.fill.fore_color.rgb = RGBColor(*bg_rgb)

    # Ensure slide_width and slide_height are valid numbers
    slide_width = prs.slide_width if prs.slide_width is not None else Inches(13.33)
    slide_height = prs.slide_height if prs.slide_height is not None else Inches(7.5)
    title_box = title_slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(float(slide_width) / Inches(1) - 2), Inches(2))
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = title_tf.paragraphs[0]
    p.text = slides[0]['title']
    p.font.bold = True
    p.font.name = font_name
    p.font.color.rgb = font_color
    p.alignment = PP_ALIGN.RIGHT if rtl else PP_ALIGN.LEFT

    title_len = len(slides[0]['title'])
    if title_len > 100:
        p.font.size = Pt(24)
    elif title_len > 60:
        p.font.size = Pt(28)
    else:
        p.font.size = Pt(36)

    if slides[0]['content']:
        subtitle = slides[0]['content'][0]
        subtitle_box = title_slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(float(slide_width) / Inches(1) - 2), Inches(1.2))
        sub_tf = subtitle_box.text_frame
        sub_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        sub_tf.word_wrap = True
        sub_tf.text = subtitle[:250]
        sp = sub_tf.paragraphs[0]
        sp.font.size = Pt(22)
        sp.font.color.rgb = font_color
        sp.font.name = font_name
        sp.alignment = PP_ALIGN.RIGHT if rtl else PP_ALIGN.LEFT

    if yt_url:
     yt_box = title_slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(float(slide_width) / Inches(1) - 2), Inches(0.5))
     yt_tf = yt_box.text_frame
     yt_tf.clear()
     p = yt_tf.paragraphs[0]
     run = p.add_run()
     run.text = "🎥 Watch related video"
     run.hyperlink.address = yt_url
     run.font.size = Pt(16)
     run.font.color.rgb = font_color
     run.font.name = font_name
     p.alignment = PP_ALIGN.RIGHT if rtl else PP_ALIGN.LEFT


    for slide in slides[1:]:
        ppt_slide = prs.slides.add_slide(prs.slide_layouts[6])
        ppt_slide.background.fill.solid()
        ppt_slide.background.fill.fore_color.rgb = RGBColor(*bg_rgb)

        # Use safe slide_width and slide_height
        image_width = Inches(float(slide_width) / Inches(1) * 0.35)
        text_width = Inches(float(slide_width) / Inches(1) * 0.60)
        margin = Inches(0.4)

        title_box = ppt_slide.shapes.add_textbox(margin, Inches(0.3), text_width - margin, Inches(1))
        title_tf = title_box.text_frame
        title_tf.clear()
        p = title_tf.paragraphs[0]
        p.text = clean_title(slide['title'])
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.underline = True
        p.font.name = font_name
        p.font.color.rgb = font_color
        p.alignment = PP_ALIGN.RIGHT if rtl else PP_ALIGN.LEFT

        content_box = ppt_slide.shapes.add_textbox(margin, Inches(1.2), text_width - margin, Inches(float(slide_height) / Inches(1) - 1.7))
        content_tf = content_box.text_frame
        content_tf.word_wrap = True
        content_tf.margin_top = Inches(0.1)
        content_tf.margin_bottom = Inches(0.1)
        content_tf.margin_left = Inches(0.1)
        content_tf.margin_right = Inches(0.1)

        points = []
        for raw_line in slide['content']:
            points.extend(split_into_bullets(raw_line))

        point_count = len(points)
        font_size = Pt(26) if point_count <= 4 else Pt(22) if point_count <= 6 else Pt(18)

        for bullet in points[:6]:
            para = content_tf.add_paragraph()
            para.text = bullet
            para.font.size = font_size
            para.font.color.rgb = font_color
            para.font.name = font_name
            para.alignment = PP_ALIGN.RIGHT if rtl else PP_ALIGN.LEFT
            para.line_spacing = Pt(font_size.pt + 6)

        img_stream = fetch_unsplash_image(subject, slide['title'], language)
        if img_stream:
            try:
                img_path = f"temp_img_{slide['title'].replace('*', '').replace(' ', '_')}.jpg"
                with open(img_path, "wb") as f:
                    f.write(img_stream.read())
                ppt_slide.shapes.add_picture(
                    img_path,
                    left=Inches(float(slide_width) / Inches(1) - float(image_width) / Inches(1)),
                    top=Inches(1.2),
                    width=Inches(float(image_width) / Inches(1) - 0.4),
                    height=Inches(float(slide_height) / Inches(1) - 1.7)
                )
                os.remove(img_path)
            except Exception as e:
                print(f"\u26a0\ufe0f Could not add image: {e}")

        footer = ppt_slide.shapes.add_textbox(margin, Inches(float(slide_height) / Inches(1) - 0.5), Inches(float(slide_width) / Inches(1)), Inches(0.5))
        footer_tf = footer.text_frame
        footer_text = {
            "english": ".",
            "hindi": ".",
            "urdu": ".",
            "marathi": ".",
        }.get(language_key, "Generated by AI Slide Generator")
        footer_tf.text = footer_text
        footer_tf.paragraphs[0].font.size = Pt(12)
        footer_tf.paragraphs[0].font.color.rgb = font_color
        footer_tf.paragraphs[0].font.name = font_name
        footer_tf.paragraphs[0].alignment = PP_ALIGN.RIGHT if rtl else PP_ALIGN.LEFT

    filename = sanitize_filename(f"Class{class_level}_{subject}_{topic}_presentation.pptx".replace(" ", ""))
    prs.save(filename)
    print(f"\u2705 PowerPoint saved as: {filename}")


def run_generator():
    print("\nAI Slide Generator")
    class_level = input("Grade level (e.g. 5): ")
    subject = input("Subject: ")
    topic = input("Topic: ")
    language = input("Language [default English]: ") or "English"
    try:
        num_slides = int(input("Number of slides [3-10, default 5]: ") or 5)
        if num_slides < 3 or num_slides > 10:
            raise ValueError
    except ValueError:
        print("\u274c Number of slides must be between 3 and 10.")
        return

    content = generate_ppt_content(class_level, subject, topic, language, num_slides)
    if isinstance(content, list):
        create_powerpoint(content, topic, class_level, subject, language)
    else:
        print(content)


if __name__ == "__main__":
    run_generator()